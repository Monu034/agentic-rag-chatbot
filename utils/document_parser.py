import pandas as pd
import PyPDF2
import docx
from pptx import Presentation
import io

def parse_pdf(file_bytes: bytes) -> str:
    text = ""
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def parse_docx(file_bytes: bytes) -> str:
    text = ""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX: {e}")
    return text

def parse_pptx(file_bytes: bytes) -> str:
    text = ""
    try:
        ppt = Presentation(io.BytesIO(file_bytes))
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text

def parse_csv(file_bytes: bytes) -> str:
    text = ""
    try:
        df = pd.read_csv(io.BytesIO(file_bytes))
        text = df.to_string()
    except Exception as e:
        print(f"Error reading CSV: {e}")
    return text

def parse_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode('utf-8')
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return ""

def parse_document(file_name: str, file_bytes: bytes) -> str:
    ext = file_name.split('.')[-1].lower()
    if ext == 'pdf':
        return parse_pdf(file_bytes)
    elif ext == 'docx':
        return parse_docx(file_bytes)
    elif ext == 'pptx':
        return parse_pptx(file_bytes)
    elif ext == 'csv':
        return parse_csv(file_bytes)
    elif ext in ['txt', 'md']:
        return parse_txt(file_bytes)
    else:
        return ""
